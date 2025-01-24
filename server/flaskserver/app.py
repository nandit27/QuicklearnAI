from flask import Flask, request, jsonify
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import pdfplumber
import re
import os
import json
from langchain_groq import ChatGroq
import google.generativeai as genai
from google.generativeai import GenerativeModel
from flask import Flask, request, jsonify
import jwt
from pymongo import MongoClient
from bson.objectid import ObjectId
from langchain_groq import ChatGroq
import os
from functools import wraps

from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from sentence_transformers import SentenceTransformer
import pdfplumber
import os
import faiss
import numpy as np
import re
import pptx
from werkzeug.utils import secure_filename

app = Flask(__name__)
SECRET_KEY = "quick" 
mongo_client = MongoClient("mongodb://localhost:27017/")  # Replace with your MongoDB URI
db = mongo_client["quicklearnai"]
topics_collection = db["statistics"]


# YouTube Transcript API configurations
formatter = TextFormatter()

def get_and_enhance_transcript(youtube_url):
    try:
        video_id = youtube_url.split('v=')[-1]
        transcript = None
        language = None

        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['hi'])
            language = 'hi'
        except:
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en'])
                language = 'en'
            except:
                return None, None

        formatted_transcript = formatter.format_transcript(transcript)

        prompt = f"""
        Act as a transcript cleaner. Generate a new transcript with the same context and the content only covered in the given transcript. 
        If there is a revision portion, differentiate it with the actual transcript.
        Give the results in sentences line by line, not in a single line. Also check whether the transcript words have any educational content relevance or not; if not then just give output as: 'Fake transcript'.
        Transcript: {formatted_transcript}
        """

        llm = ChatGroq(
            model="llama-3.1-70b-versatile",
            temperature=0,
            groq_api_key="YOUR_GROQ_API_KEY"
        )

        enhanced_transcript = llm.invoke(prompt)

        return enhanced_transcript, language
    except Exception as e:
        print(f"Error: {str(e)}")
        return None, None

def generate_summary_and_quiz(transcript, num_questions, language, difficulty):
    try:
        prompt = f"""
        Summarize the following transcript by identifying the key topics covered, and provide a detailed summary of each topic in 6-7 sentences.
        Each topic should be labeled clearly as "Topic X", where X is the topic name. Provide the full summary for each topic in English, even if the transcript is in a different language.
        Strictly ensure that possessives (e.g., John's book) and contractions (e.g., don't) use apostrophes (`'`) instead of quotation marks (" or “ ”).

        If the transcript contains 'Fake Transcript', do not generate any quiz or summary.

        After the summary, give the name of the topic on which the transcript was all about in a maximum of 2 to 3 words.
        After summarizing, create a quiz with {num_questions} multiple-choice questions in English, based on the transcript content.
        Only generate {difficulty} difficulty questions. Format the output in JSON format as follows, just give the JSON as output, nothing before it:

        {{
            "summary": {{
                "topic1": "value1",
                "topic2": "value2",
                "topic3": "value3"
            }},
            "questions": {{
                "{difficulty}": [
                    {{
                        "question": "What is the capital of France?",
                        "options": ["Paris", "London", "Berlin", "Madrid"],
                        "answer": "Paris"
                    }},
                    {{
                        "question": "What is the capital of Germany?",
                        "options": ["Paris", "London", "Berlin", "Madrid"],
                        "answer": "Berlin"
                    }}
                ]
            }}
        }}

        Transcript: {transcript}
        """

        llm = ChatGroq(
            model="llama-3.1-70b-versatile",
            temperature=0,
            groq_api_key="YOUR_GROQ_API_KEY"
        )
        response = llm.invoke(prompt)

        if hasattr(response, 'content'):
            response_content = response.content
            try:
                response_dict = json.loads(response_content)
                return response_dict
            except json.JSONDecodeError as e:
                print(f"JSONDecodeError: {e}")
                return None
        else:
            print("Response does not have a 'content' attribute.")
            return None

    except Exception as e:
        print(f"Error generating summary and quiz: {str(e)}")
        return None

# PDF Q&A configurations
from flask import Flask, request, jsonify
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from sentence_transformers import SentenceTransformer
import pdfplumber
import os
import faiss
import numpy as np
import re
import pptx

app = Flask(__name__)

# Models and Globals
embedding_model = SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')
tokenizer = AutoTokenizer.from_pretrained("t5-small")
generator_model = AutoModelForSeq2SeqLM.from_pretrained("t5-small")

index = None
corpus = []

# PDF Processing
def load_pdf_and_split(file_path, chunk_size=200):
    full_text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    full_text += text + " "
    except Exception as e:
        print(f"Error reading PDF: {str(e)}")
        return []

    full_text = re.sub(r'\s+', ' ', full_text).strip()
    chunks = []
    sentences = re.split(r'(?<=[.!?])\s+', full_text)
    current_chunk = ""

    for sentence in sentences:
        if len(current_chunk) + len(sentence) < chunk_size:
            current_chunk += sentence + " "
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = sentence + " "
    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks


def load_ppt_and_split(file_path, chunk_size=200):
    full_text = ""
    try:
        presentation = pptx.Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text += shape.text.strip() + " "
    except Exception as e:
        print(f"Error reading PPT: {str(e)}")
        return []

    full_text = re.sub(r'\s+', ' ', full_text).strip()
    return load_pdf_and_split(full_text, chunk_size)


def build_faiss_index(corpus):
    corpus_embeddings = embedding_model.encode(corpus)
    dimension = corpus_embeddings.shape[1]
    faiss_index = faiss.IndexFlatL2(dimension)
    faiss_index.add(np.array(corpus_embeddings).astype('float32'))
    return faiss_index, corpus_embeddings

# File Upload Endpoint
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}

def is_allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/upload', methods=['POST'])
def upload_file():
    global index, corpus

    print(request.files)
        
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request."}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected for uploading."}), 400

    if not is_allowed_file(file.filename):
        return jsonify({"error": "Invalid file type. Please upload a PDF or PPTX file."}), 400

    try:
        # Save the file securely
        filename = secure_filename(file.filename)
        file_path = os.path.join("uploads", filename)
        os.makedirs("uploads", exist_ok=True)
        file.save(file_path)

        # Process the file
        if filename.endswith('.pdf'):
            corpus = load_pdf_and_split(file_path)
        else:
            corpus = load_ppt_and_split(file_path)

        if not corpus:
            return jsonify({"error": "Could not extract text from the file."}), 500

        # Build the FAISS index
        index, _ = build_faiss_index(corpus)
        return jsonify({"message": "File processed and index built successfully."}), 200

    except Exception as e:
        return jsonify({"error": f"An error occurred while processing the file: {str(e)}"}), 500


@app.route('/ask', methods=['POST'])
def ask_question():
    global index, corpus
    if index is None or not corpus:
        return jsonify({"error": "No file uploaded or index built yet."}), 400

    query = request.json.get('query')
    if not query:
        return jsonify({"error": "No query provided."}), 400

    try:
        query_embedding = embedding_model.encode([query])
        distances, indices = index.search(np.array(query_embedding).astype('float32'), k=3)
        relevant_chunks = [corpus[i] for i in indices[0] if i < len(corpus)]

        input_text = " ".join(relevant_chunks) + " Question: " + query
        input_ids = tokenizer.encode(input_text, return_tensors='pt', truncation=True, max_length=512)
        output_ids = generator_model.generate(input_ids, max_length=150, num_beams=4, early_stopping=True)
        answer = tokenizer.decode(output_ids[0], skip_special_tokens=True)

        return jsonify({"answer": answer}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500




@app.route('/quiz', methods=['POST'])
def generate_quiz():
    data = request.json
    youtube_link = data.get('youtube_url')
    num_questions = data.get('num_questions', 5)
    difficulty = data.get('difficulty', 'medium')

    if youtube_link:
        transcript, language = get_and_enhance_transcript(youtube_link)
        
        if transcript:
            summary_and_quiz = generate_summary_and_quiz(transcript, num_questions, language, difficulty)
            if summary_and_quiz:
                return jsonify(summary_and_quiz)
            else:
                return jsonify({"error": "Failed to generate quiz"}), 500
        else:
            return jsonify({"error": "Failed to fetch transcript"}), 404
    else:
        return jsonify({"error": "No YouTube URL provided"}), 400
 # Replace with your actual secret key


# Middleware for token validation
def validate_token_middleware():
    def middleware(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            auth_header = request.headers.get("Authorization")
            token = auth_header.split("Bearer ")[-1] if auth_header and "Bearer " in auth_header else None
            
            if not token:
                return jsonify({"message": "Unauthorized: No token provided"}), 401
            
            try:
                decoded = jwt.decode(token, SECRET_KEY, algorithms=["HS256"])
                request.user_id = decoded.get("id")
                request.user_role = decoded.get("role")  # Optional
                return func(*args, **kwargs)
            except jwt.ExpiredSignatureError:
                return jsonify({"message": "Unauthorized: Token has expired"}), 401
            except jwt.InvalidTokenError as e:
                print(f"Token decoding error: {e}")
                return jsonify({"message": "Unauthorized: Invalid token"}), 401
        
        return wrapper
    return middleware


# Function to interact with LLaMA API
def llama_generate_recommendations(prompt):
    try:
        # Configure the API key
        genai.configure(api_key=os.getenv("GENAI_API_KEY"))
        
        # Create Gemini Flash model instance
        model = GenerativeModel('gemini-2.0-flash-exp')
        
        # Generate response
        response = model.generate_content(prompt)
        
        return response.text
    except Exception as e:
        return f"Error connecting to Gemini API: {e}"
    
@app.route('/getonly', methods=['GET'])
@validate_token_middleware()
def get_recommendations():
    user_id = request.user_id  # Extract user ID from the token
    try:
        user_documents = topics_collection.find({"student": ObjectId(user_id)})
        user_list = list(user_documents)

        topics = [doc.get("topic") for doc in user_list if "topic" in doc]

        if not topics:
            return jsonify({"message": "No topics found for the provided user."}), 404

        prompt = f"Act as a recommendation generator , generate and recommend content for the following topics , also give five urls of YouTube videos regarding the topic .If there are multiple topics give overview of each of them and links for each topic video as well. The topics are: {', '.join(topics)}"
        recommendations = llama_generate_recommendations(prompt)

        return jsonify({
            "message": "Recommendations generated successfully",
            "recommendations": recommendations
        }), 200

    except Exception as e:
        print("Error:", str(e))
        return jsonify({"message": f"An error occurred: {str(e)}"}), 500


@app.route('/')
def home():
    return jsonify({"message": "Welcome to the QuickLearn AI server"}), 200

if __name__ == '__main__':
    app.run(debug=True)
